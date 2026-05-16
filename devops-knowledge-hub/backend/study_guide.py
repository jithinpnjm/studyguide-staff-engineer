"""
Study guide generator — reads every PDF, identifies concept gaps,
searches the web to fill them, then synthesises one definitive "zero to hero" guide per domain.
"""

import asyncio
import json
import logging
from pathlib import Path

import ai_client
import database

log = logging.getLogger(__name__)

STUDY_GUIDES_DIR = Path(__file__).parent.parent / "study-guides"

DOMAIN_LABELS = {
    "kubernetes-containers":   "Kubernetes & Containers",
    "cicd-gitops":             "CI/CD & GitOps",
    "observability":           "Observability",
    "cloud-platforms":         "Cloud Platforms",
    "infrastructure-as-code":  "Infrastructure as Code",
    "security-devsecops":      "Security & DevSecOps",
    "networking-service-mesh": "Networking & Service Mesh",
    "sre-practices":           "SRE Practices",
    "linux-systems":           "Linux & Systems",
    "databases-storage":       "Databases & Storage",
    "platform-engineering":    "Platform Engineering",
    "ai-mlops":                "AI & MLOps",
    "agentic-genai":           "Agentic & GenAI",
    "llmops":                  "LLMOps",
    "scripting-python":        "Python for DevOps",
    "scripting-bash-shell":    "Bash & Shell Scripting",
    "scripting-go":            "Go for DevOps",
    "automation-tooling":      "Automation & Tooling",
    "general-devops":          "General DevOps",
}

DOMAIN_EXTRAS = {
    "kubernetes-containers":   "Pod scheduling, resource management, operators, multi-tenancy, GitOps delivery, cluster reliability at scale, HPA/VPA/KEDA, network policies, admission controllers.",
    "cicd-gitops":             "Pipeline design, GitOps reconciliation loops (Flux/ArgoCD), progressive delivery, secrets management in CI, supply chain security (SLSA/Sigstore).",
    "observability":           "Three pillars (metrics/logs/traces), SLO-based alerting, cardinality management, distributed tracing, sampling strategies, observability-driven debugging.",
    "cloud-platforms":         "AWS/GCP/Azure from an SRE lens: managed services trade-offs, landing zones, cost optimisation, IAM at scale, multi-cloud patterns, egress costs.",
    "infrastructure-as-code":  "Terraform, Pulumi, CDK: state management, module design, drift detection, testing with Terratest, IaC at enterprise scale, GitOps for infra.",
    "security-devsecops":      "RBAC, OPA/Kyverno policy-as-code, supply chain security, secrets management (Vault/ESO), zero-trust, SBOM, container image signing.",
    "networking-service-mesh": "CNI plugins, eBPF networking, Envoy proxy internals, Istio/Linkerd, mTLS, traffic shaping, east-west networking, Gateway API.",
    "sre-practices":           "SLO/SLI/error budgets, toil elimination, incident management, blameless post-mortems, capacity planning, chaos engineering, runbooks.",
    "linux-systems":           "Kernel internals for containers (cgroups v2, namespaces), performance tuning (perf/eBPF), systemd, IO schedulers, memory subsystem.",
    "databases-storage":       "Database operations, replication topologies, sharding strategies, backup/recovery, CSI drivers, storage classes, stateful workloads on K8s.",
    "platform-engineering":    "Internal developer platforms, golden paths, self-service infra, platform SLOs, Backstage, Crossplane, and staff engineer platform leadership.",
    "ai-mlops":                "ML pipeline automation (Kubeflow/MLflow), model serving (Triton/TorchServe), GPU scheduling on K8s, feature stores, model observability.",
    "agentic-genai":           "Agent architectures (ReAct/CoT), tool use, memory/retrieval patterns, multi-agent orchestration, production deployment, evaluation.",
    "llmops":                  "LLM inference optimisation (vLLM/TensorRT), RAG pipeline design, prompt versioning, evaluation frameworks, cost/latency trade-offs.",
    "scripting-python":        "Python for infra automation: Boto3/SDK patterns, async programming, CLI tooling with Click/Typer, testing infra code with pytest.",
    "scripting-bash-shell":    "Shell scripting best practices, POSIX portability, error handling with set -euo pipefail, signal trapping, parallel execution, common DevOps patterns.",
    "scripting-go":            "Go for DevOps: CLI tools with Cobra/Viper, Kubernetes controllers with controller-runtime, concurrency patterns, building reliable daemons.",
    "automation-tooling":      "Ansible at scale, event-driven automation (EDA), Crossplane, configuration management trade-offs, idempotency, and automation architecture.",
    "general-devops":          "DevOps transformation, DORA metrics, team topologies, platform teams, psychological safety, and staff engineer influence without authority.",
}

GAP_PROMPT = """You are a Staff/Principal SRE/DevOps engineer reviewing source material for the domain: {domain_label}.

Here is what the existing documents ALREADY cover:
{existing_topics}

List the TOP 5 most important concepts/topics that are MISSING or UNDER-COVERED for a complete
staff-level mastery of {domain_label}. These are gaps we need to fill from the web.

Return ONLY a JSON array of search queries, e.g.:
["kubernetes admission controllers production patterns", "kube-scheduler internals staff level"]
"""

SYNTHESIS_PROMPT = """You are synthesising the definitive "Zero to Hero" study guide for a Staff/Principal SRE/DevOps engineer.

Domain: {domain_label}
{domain_extra}

You have TWO sources of material:

=== SOURCE A: Content from {doc_count} personal documents ===
{doc_context}

=== SOURCE B: Web research to fill gaps ===
{web_context}

Synthesise EVERYTHING above into ONE authoritative, comprehensive guide.
Rules:
- DEDUPLICATE: if the same concept appears in multiple sources, merge it into one clear explanation
- FILL GAPS: use the web research to cover topics missing from personal docs
- DEPTH over breadth: go deep on each concept rather than listing superficially
- NO installation steps, no "how to install X" content
- Write for a Staff/Principal engineer — assume strong baseline knowledge
- Use concrete examples, production patterns, failure scenarios

Write using EXACTLY this structure (all sections mandatory, all must be detailed):

# {domain_label} — Zero to Hero

## 🎯 Why This Domain Matters
What problems does mastering this solve? What business/reliability outcomes depend on it?
Why do staff engineers need deep expertise here? What goes wrong without it?

## 📋 Prerequisites & Mental Models
What conceptual foundation unlocks this domain? Key abstractions to internalise first.
How to think about this domain (the mental model that makes everything else click).

## 🔷 Core Concepts
Every fundamental abstraction, primitive, and principle — explained with depth.
Not definitions — explain WHY each concept exists and its real-world implications.

## 🛠️ Tools & Ecosystem
Primary tools and frameworks. For each: what it solves, when to use it, when NOT to use it,
key trade-offs, and cloud-managed vs self-hosted decisions.

## 🏗️ Architecture Patterns
Production-proven patterns. For each: diagram description, trade-offs, when to choose it,
what breaks it at scale, migration paths.

## ⚙️ Production Operations & Day-2
What does running this in production actually look like?
Capacity planning, scaling triggers, rollout strategies, configuration tuning, runbook patterns.

## 📊 Observability & Debugging
Critical metrics, golden signals, alerting strategy.
Step-by-step debugging workflow for the top 5 failure scenarios in this domain.

## 🔐 Security Considerations
Threat model, attack surfaces, hardening checklist, common misconfigurations, compliance hooks.

## 🎓 Staff/Principal Engineer Perspective
Platform-level thinking. Architectural decision frameworks. How to lead this domain.
What decisions belong to staff vs senior engineers. How to evaluate options at scale.

## 💥 Failure Modes & Incident Patterns
How this breaks in production — specific failure modes with: early signals, blast radius,
debugging steps, and prevention. Real patterns from post-mortems.

## 💼 Interview & Design Review Prep
System design questions that probe this domain. Expected staff-level depth of answer.
Common follow-up questions. Trade-off discussions. What interviewers are testing.

## 📚 Key Takeaways
The 15 most important principles. Written as actionable wisdom, not bullet facts.
"""


def _extract_file_text(file_path: str, file_type: str, max_chars: int = 2500) -> str:
    path = Path(file_path)
    if not path.exists():
        return ""
    try:
        ft = (file_type or "").lower()
        if ft == "pdf":
            import fitz
            doc = fitz.open(str(path))
            text = ""
            for page in doc[:12]:
                text += page.get_text()
                if len(text) >= max_chars:
                    break
            doc.close()
            return text[:max_chars]
        if ft in ("pptx", "ppt"):
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = [
                shape.text.strip()
                for slide in list(prs.slides)[:20]
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            return "\n".join(parts)[:max_chars]
        if ft in ("docx", "doc"):
            from docx import Document
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())[:max_chars]
        if ft in ("md", "markdown", "txt"):
            return path.read_text(errors="replace")[:max_chars]
    except Exception as e:
        log.debug("Text extraction failed for %s: %s", file_path, e)
    return ""


def _build_doc_context(domain: str, max_total: int = 100000) -> tuple[str, int, list[str]]:
    """Extract content from all docs in domain. Returns (context, doc_count, topics_covered)."""
    docs = database.get_docs_by_domain(domain, doc_type="raw")
    if not docs:
        return "", 0, []

    parts = []
    used = 0
    all_topics: list[str] = []
    per_doc = min(2500, max_total // max(len(docs), 1))

    for doc in docs:
        title = doc.get("title") or doc.get("filename") or "Untitled"

        # Collect topics for gap analysis
        for t in (doc.get("key_topics") or []):
            if t not in all_topics:
                all_topics.append(t)

        # Best content source: concept card → AI summary → raw file text
        card = database.get_concept_card_for_doc(doc["id"])
        if card and card.get("raw_text"):
            text = card["raw_text"][:per_doc]
        elif doc.get("summary"):
            topics = ", ".join(doc.get("key_topics") or [])
            text = f"{doc['summary']}\nTopics: {topics}"[:per_doc]
        elif doc.get("file_path"):
            text = _extract_file_text(doc["file_path"], doc.get("file_type", ""), per_doc)
        else:
            continue

        if not text.strip():
            continue

        chunk = f"--- [{title}] ---\n{text}\n"
        if used + len(chunk) > max_total:
            break
        parts.append(chunk)
        used += len(chunk)

    return "\n".join(parts), len(docs), all_topics


async def _fetch_web_content(url: str, max_chars: int = 3000) -> str:
    """Fetch and extract clean text from a URL."""
    try:
        import httpx
        async with httpx.AsyncClient(timeout=10, follow_redirects=True,
                                     headers={"User-Agent": "Mozilla/5.0"}) as client:
            resp = await client.get(url)
            if resp.status_code != 200:
                return ""
            html = resp.text
        try:
            import trafilatura
            text = trafilatura.extract(html, include_comments=False, include_tables=True)
            return (text or "")[:max_chars]
        except Exception:
            pass
        # Fallback: strip HTML tags
        import re
        text = re.sub(r"<[^>]+>", " ", html)
        text = re.sub(r"\s+", " ", text).strip()
        return text[:max_chars]
    except Exception as e:
        log.debug("Failed to fetch %s: %s", url, e)
        return ""


async def _search_and_fetch(query: str, num_results: int = 3) -> str:
    """DuckDuckGo search + fetch top results → combined text."""
    try:
        from ddgs import DDGS
        loop = asyncio.get_event_loop()

        def _search():
            return list(DDGS().text(query, max_results=num_results))

        hits = await loop.run_in_executor(None, _search)
        if not hits:
            return ""

        fetch_tasks = [_fetch_web_content(h["href"]) for h in hits if h.get("href")]
        contents = await asyncio.gather(*fetch_tasks, return_exceptions=True)

        parts = []
        for hit, content in zip(hits, contents):
            if isinstance(content, str) and content.strip():
                parts.append(f"[{hit.get('title', '')}]\n{content}")

        return "\n\n".join(parts)
    except Exception as e:
        log.warning("Web search failed for '%s': %s", query, e)
        return ""


async def _find_and_fill_gaps(domain_label: str, existing_topics: list[str]) -> str:
    """Ask Gemini what's missing, search for it, return enriched web content."""
    if not existing_topics:
        # No topics yet — search for core curriculum
        queries = [
            f"{domain_label} fundamentals staff engineer",
            f"{domain_label} production architecture patterns",
            f"{domain_label} failure modes incident patterns",
        ]
    else:
        # Ask Gemini to identify gaps
        gap_prompt = GAP_PROMPT.format(
            domain_label=domain_label,
            existing_topics=", ".join(existing_topics[:40]),
        )
        try:
            raw = await ai_client.generate(gap_prompt, task="answer_question", json_output=True)
            import re
            match = re.search(r'\[.*\]', raw, re.DOTALL)
            queries = json.loads(match.group()) if match else []
        except Exception:
            queries = []

        # Always add these critical queries regardless of gap analysis
        queries.extend([
            f"{domain_label} staff principal engineer deep dive",
            f"{domain_label} production failure modes",
        ])
        queries = queries[:5]

    log.info("Web gap searches for %s: %s", domain_label, queries)
    search_tasks = [_search_and_fetch(q, num_results=2) for q in queries]
    results = await asyncio.gather(*search_tasks, return_exceptions=True)

    web_parts = []
    for query, content in zip(queries, results):
        if isinstance(content, str) and content.strip():
            web_parts.append(f"=== Web: {query} ===\n{content}")

    return "\n\n".join(web_parts)


async def generate_study_guide(domain: str, min_docs: int = 1) -> dict | None:
    if not ai_client.is_configured():
        log.error("Gemini not configured — cannot generate study guide")
        return None

    docs = database.get_docs_by_domain(domain, doc_type="raw")
    if len(docs) < min_docs:
        return None

    domain_label = DOMAIN_LABELS.get(domain, domain)
    domain_extra = DOMAIN_EXTRAS.get(domain, "")

    log.info("Building context for %s from %d documents…", domain_label, len(docs))
    doc_context, doc_count, existing_topics = _build_doc_context(domain)

    log.info("Searching web to fill concept gaps for %s…", domain_label)
    web_context = await _find_and_fill_gaps(domain_label, existing_topics)

    if not doc_context and not web_context:
        log.warning("No content available for %s — skipping", domain)
        return None

    prompt = SYNTHESIS_PROMPT.format(
        domain_label=domain_label,
        domain_extra=domain_extra,
        doc_count=doc_count,
        doc_context=doc_context or "(No document summaries available yet — run categorization first)",
        web_context=web_context or "(No web content retrieved)",
    )

    log.info("Synthesising guide for %s (%d docs, %d chars context)…",
             domain_label, doc_count, len(doc_context) + len(web_context))
    try:
        content = await ai_client.generate(prompt, task="study_guide")
    except Exception as e:
        log.error("Guide generation failed for %s: %s", domain, e)
        return None

    STUDY_GUIDES_DIR.mkdir(parents=True, exist_ok=True)
    (STUDY_GUIDES_DIR / f"{domain}.md").write_text(content, encoding="utf-8")

    database.save_study_guide(domain, content, doc_count)
    log.info("✓ %s — %d chars written", domain_label, len(content))
    return {"domain": domain, "label": domain_label, "doc_count": doc_count, "chars": len(content)}


async def generate_all_study_guides(min_docs: int = 1, status: dict | None = None) -> list[dict]:
    domains_with_docs = [d for d in DOMAIN_LABELS if database.get_docs_by_domain(d, doc_type="raw")]
    if status is not None:
        status["total"] = len(domains_with_docs)

    results = []
    for domain in DOMAIN_LABELS:
        docs = database.get_docs_by_domain(domain, doc_type="raw")
        if not docs:
            continue
        if status is not None:
            status["current"] = DOMAIN_LABELS.get(domain, domain)
        result = await generate_study_guide(domain, min_docs=min_docs)
        if result:
            results.append(result)
        if status is not None:
            status["done"] = len(results)
        await asyncio.sleep(5)
    return results


NEW_DOC_PATCH_PROMPT = """An existing "Zero to Hero" study guide already covers {domain_label} in depth.

The guide currently contains these sections:
{existing_sections}

A NEW document has just been added to this domain:
Title: {doc_title}
Content extracted from the document:
---
{doc_content}
---

Your task: Write ONLY what is GENUINELY NEW — concepts, patterns, tools, or insights that are
NOT already covered in the existing guide sections above.

Rules:
- If the document only repeats existing content → respond with exactly: NO_NEW_CONTENT
- Do NOT rewrite or summarise existing sections
- Do NOT duplicate anything already in the guide
- Focus only on what this specific document adds that is missing
- Write at staff/principal engineer depth
- Be specific: what new mental model, pattern, or production insight does this add?

Format your response as:

## 🆕 New Learnings: {doc_title}

> *Added {date} — what this document adds beyond the existing guide*

[Your content here — only genuinely new concepts/patterns/insights]
"""


async def patch_study_guide_with_doc(domain: str, doc_id: str) -> dict | None:
    """
    Smart append: reads a new doc, compares to existing guide,
    appends ONLY genuinely new sections. User doesn't need to re-read the whole guide.
    """
    if not ai_client.is_configured():
        return None

    existing_guide = database.get_study_guide(domain)
    if not existing_guide:
        # No guide yet — generate fresh
        return await generate_study_guide(domain, min_docs=1)

    doc = database.get_document(doc_id)
    if not doc:
        return None

    domain_label = DOMAIN_LABELS.get(domain, domain)
    doc_title = doc.get("title") or doc.get("filename") or "New Document"

    # Extract content from the new document
    card = database.get_concept_card_for_doc(doc_id)
    if card and card.get("raw_text"):
        doc_content = card["raw_text"][:4000]
    elif doc.get("summary"):
        topics = ", ".join(doc.get("key_topics") or [])
        doc_content = f"{doc['summary']}\nKey topics: {topics}"
    elif doc.get("file_path"):
        doc_content = _extract_file_text(doc["file_path"], doc.get("file_type", ""), 4000)
    else:
        return None

    if not doc_content.strip():
        return None

    # Collect existing section headings
    existing_sections = [
        line.lstrip("# ").strip()
        for line in existing_guide["content"].split("\n")
        if line.startswith("## ") or line.startswith("### ")
    ]

    from datetime import date
    prompt = NEW_DOC_PATCH_PROMPT.format(
        domain_label=domain_label,
        existing_sections="\n".join(f"- {s}" for s in existing_sections),
        doc_title=doc_title,
        doc_content=doc_content,
        date=date.today().strftime("%Y-%m-%d"),
    )

    try:
        new_section = await ai_client.generate(prompt, task="study_guide")
    except Exception as e:
        log.error("Patch generation failed: %s", e)
        return None

    if "NO_NEW_CONTENT" in new_section:
        log.info("New doc '%s' adds nothing new to %s guide — skipping patch", doc_title, domain)
        return {"domain": domain, "patched": False, "reason": "no_new_content", "doc_title": doc_title}

    # Append the new section
    updated_content = existing_guide["content"].rstrip() + f"\n\n---\n\n{new_section.strip()}\n"

    STUDY_GUIDES_DIR.mkdir(parents=True, exist_ok=True)
    (STUDY_GUIDES_DIR / f"{domain}.md").write_text(updated_content, encoding="utf-8")

    database.save_study_guide(domain, updated_content, existing_guide.get("doc_count", 0) + 1)
    database.save_study_guide_patch(domain, doc_id, "new_section", new_section)

    log.info("✓ Appended new section to %s guide for: %s", domain, doc_title)
    return {
        "domain": domain,
        "patched": True,
        "doc_title": doc_title,
        "new_section_preview": new_section[:400],
    }
