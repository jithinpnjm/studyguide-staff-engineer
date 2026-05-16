"""
Phase 2 — AI Categorization Engine (Gemini).
Assigns domain, tags, summary, concept_level, key_topics.
"""

import asyncio
import json
import logging
import os
from datetime import datetime, timezone
from pathlib import Path

import ai_client
import database

log = logging.getLogger(__name__)

DOMAINS = [
    "kubernetes-containers", "cicd-gitops", "observability", "cloud-platforms",
    "infrastructure-as-code", "security-devsecops", "networking-service-mesh",
    "sre-practices", "linux-systems", "databases-storage", "platform-engineering",
    "ai-mlops", "agentic-genai", "llmops",
    "scripting-python", "scripting-bash-shell", "scripting-go",
    "automation-tooling", "general-devops",
]

TIER = os.getenv("GEMINI_TIER", "free")
BATCH_CONFIG = {
    "free": {"batch_size": 10, "delay_between": 4.0, "delay_batch": 15.0},
    "paid": {"batch_size": 25, "delay_between": 0.5, "delay_batch": 2.0},
}

SYSTEM_CONTEXT = (
    "You are an expert SRE/DevOps knowledge classifier for a Senior/Staff/Principal "
    "DevOps-SRE-Cloud Engineer's personal study portal. Classify documents with a "
    "staff-engineer lens — focus on concepts, architecture, trade-offs, and production "
    "relevance. Respond ONLY in valid JSON, no markdown, no explanation outside the JSON."
)

USER_PROMPT_TEMPLATE = """Analyze this document and return exactly this JSON structure:
{{
  "domain": "<one value from taxonomy>",
  "tags": ["tag1", "tag2", "tag3"],
  "summary": "<3 sentences: what this covers, key concepts taught, who benefits. Staff-engineer perspective. No installation steps. No fluff.>",
  "concept_level": "beginner|intermediate|advanced|staff-level",
  "key_topics": ["topic1", "topic2", "...up to 8 topics"],
  "title_suggestion": "<clean title if filename is unclear, else null>"
}}

Domain taxonomy (choose exactly one):
{domains}

Document filename: {filename}
Document title: {title}

Document content:
{text}"""

FINGERPRINT_PROMPT = """Extract a quick knowledge fingerprint from this document.
Return JSON only:
{{
  "probable_domain": "<one taxonomy value>",
  "core_topics": ["topic1", "...up to 10"],
  "key_technologies": ["tech1", "...up to 10"],
  "summary_one_line": "<15 words max>"
}}

Domain taxonomy:
{domains}

First 2000 characters of document:
{text}"""


# ── Text extraction ──────────────────────────────────────────────────────────

def extract_text_from_pdf(file_path: str, max_pages: int = 20) -> str:
    try:
        import fitz
        doc = fitz.open(file_path)
        pages = [doc[i].get_text() for i in range(min(max_pages, len(doc)))]
        doc.close()
        return "\n\n".join(pages)
    except Exception as e:
        log.warning("PDF extraction failed %s: %s", file_path, e)
        return ""


def extract_text_from_pptx(file_path: str, max_slides: int = 30) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides[:max_slides]:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts)
    except Exception as e:
        log.warning("PPTX extraction failed %s: %s", file_path, e)
        return ""


def extract_text_from_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        log.warning("DOCX extraction failed %s: %s", file_path, e)
        return ""


def extract_text(doc: dict) -> str:
    raw = doc.get("raw_text") or ""
    if len(raw) > 500:
        return raw[:12000]

    file_path = doc.get("file_path", "")
    file_type = doc.get("file_type", "").lower()

    if not file_path or not Path(file_path).exists():
        return raw

    if file_type == "pdf":
        return extract_text_from_pdf(file_path)[:12000]
    if file_type in ("pptx", "ppt"):
        return extract_text_from_pptx(file_path)[:12000]
    if file_type in ("docx", "doc"):
        return extract_text_from_docx(file_path)[:12000]
    if file_type in ("md", "txt", "markdown"):
        try:
            return Path(file_path).read_text(errors="replace")[:12000]
        except Exception:
            return raw
    return raw


# ── Fingerprint (cheap, fast, pre-dedup) ────────────────────────────────────

async def get_fingerprint(text: str) -> dict | None:
    if not ai_client.is_configured():
        return None
    prompt = FINGERPRINT_PROMPT.format(
        domains=", ".join(DOMAINS),
        text=text[:2000],
    )
    try:
        raw = await ai_client.generate(prompt, task="categorize", json_output=True)
        return _safe_parse(raw, "fingerprint")
    except Exception as e:
        log.debug("Fingerprint failed: %s", e)
        return None


# ── Core categorization ──────────────────────────────────────────────────────

def _build_prompt(doc: dict, text: str) -> str:
    return f"{SYSTEM_CONTEXT}\n\n" + USER_PROMPT_TEMPLATE.format(
        domains="\n".join(f"  {d}" for d in DOMAINS),
        title=doc.get("title") or "",
        filename=doc.get("filename") or "",
        text=text[:10000],
    )


def _safe_parse(raw: str, context: str = "") -> dict | None:
    raw = raw.strip()
    if raw.startswith("```"):
        parts = raw.split("```")
        raw = parts[1] if len(parts) > 1 else raw
        if raw.startswith("json"):
            raw = raw[4:]
        raw = raw.strip().rstrip("`").strip()
    try:
        return json.loads(raw)
    except json.JSONDecodeError as e:
        log.warning("JSON parse failed (%s): %s | raw=%s", context, e, raw[:200])
        return None


async def categorize_document(doc_id: str) -> dict | None:
    if not ai_client.is_configured():
        log.debug("Gemini not configured — skipping %s", doc_id)
        return None

    doc = database.get_document(doc_id)
    if not doc:
        log.warning("Document not found: %s", doc_id)
        return None

    text = extract_text(doc)
    prompt = _build_prompt(doc, text)
    file_path = doc.get("file_path", "")
    file_type = doc.get("file_type", "").lower()

    data = None

    # Try native PDF upload for image-heavy docs
    if file_type == "pdf" and file_path and Path(file_path).exists():
        try:
            raw = ai_client.categorize_pdf_native(file_path, prompt)
            data = _safe_parse(raw, f"native-pdf:{doc_id[:8]}")
            if data:
                log.debug("Native PDF upload succeeded for %s", doc_id[:8])
        except Exception as e:
            log.info("Native PDF upload failed for %s, falling back to text: %s", doc_id[:8], e)

    # Text-based categorization (fallback or non-PDF)
    if data is None:
        try:
            raw = await ai_client.generate(prompt, task="categorize", json_output=True)
            data = _safe_parse(raw, f"text:{doc_id[:8]}")
        except Exception as e:
            log.warning("Gemini API error for %s: %s", doc_id[:8], e)
            database.log_categorization_error(doc_id, str(e))
            return None

    if data is None:
        database.log_categorization_error(doc_id, "JSON parse failed after both strategies")
        return None

    domain = data.get("domain", "general-devops")
    if domain not in DOMAINS:
        log.warning("Unknown domain '%s' for %s — using general-devops", domain, doc_id[:8])
        domain = "general-devops"

    tags = [str(t) for t in (data.get("tags") or [])[:8]]
    key_topics = [str(t) for t in (data.get("key_topics") or [])[:8]]

    updates: dict = {
        "domain": domain,
        "tags": tags,
        "summary": (data.get("summary") or "")[:600] or None,
        "concept_level": data.get("concept_level") or "intermediate",
        "key_topics": key_topics,
        "categorized_at": datetime.now(timezone.utc).isoformat(),
        "dedup_status": "unique",
    }

    # Overwrite title only if it looks like a raw filename
    current_title = doc.get("title") or ""
    title_suggestion = data.get("title_suggestion")
    looks_like_filename = (
        not current_title
        or (" " not in current_title and "." in current_title)
        or current_title.strip().isdigit()
    )
    if title_suggestion and looks_like_filename:
        updates["title"] = title_suggestion

    database.update_document(doc_id, updates)
    log.info(
        "✓ %-36s  →  %-32s  [%s]",
        doc_id[:36],
        domain,
        updates["concept_level"],
    )
    return database.get_document(doc_id)


async def recategorize_document(doc_id: str) -> dict | None:
    return await categorize_document(doc_id)


async def batch_categorize(
    limit: int = 1000,
    only_uncategorized: bool = True,
    progress_callback=None,
) -> dict:
    cfg = BATCH_CONFIG.get(TIER, BATCH_CONFIG["free"])

    docs = database.get_uncategorized_docs(limit=limit) if only_uncategorized \
        else database.list_documents(limit=limit, doc_type="raw")

    total = len(docs)
    processed = 0
    failed = []

    log.info("Batch: %d documents  tier=%s  delay=%.1fs", total, TIER, cfg["delay_between"])

    for i, doc in enumerate(docs):
        doc_id = doc["id"]
        try:
            result = await categorize_document(doc_id)
            if result is None:
                failed.append(doc_id)
            else:
                processed += 1
                if progress_callback:
                    await progress_callback(
                        processed=processed,
                        total=total,
                        current_title=result.get("title") or doc.get("filename"),
                        current_domain=result.get("domain"),
                    )
        except Exception as e:
            log.error("Batch error %s: %s", doc_id[:8], e)
            database.log_categorization_error(doc_id, str(e))
            failed.append(doc_id)

        await asyncio.sleep(cfg["delay_between"])
        if (i + 1) % cfg["batch_size"] == 0 and i < total - 1:
            log.info("Batch pause at %d/%d (processed=%d)", i + 1, total, processed)
            await asyncio.sleep(cfg["delay_batch"])

    log.info("Batch done: processed=%d  failed=%d  total=%d", processed, len(failed), total)
    return {"processed": processed, "failed": failed, "total": total}
