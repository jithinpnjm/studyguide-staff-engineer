import asyncio
import base64
import hashlib
import html as _html
import json
import logging
import mimetypes
import os
import threading
import uuid
from pathlib import Path
from typing import AsyncGenerator, Optional

import fitz  # PyMuPDF
import uvicorn
from dotenv import load_dotenv
from fastapi import BackgroundTasks, FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse, Response, StreamingResponse, FileResponse

load_dotenv()

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)-8s %(name)s — %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger(__name__)

import database
from categorizer import batch_categorize, categorize_document, recategorize_document
from models import (
    BatchCategorizeRequest,
    DiscoverDownloadRequest,
    DiscoverRequest,
    DocumentListItem,
    DocumentRecord,
    DomainUpdate,
    ImageUploadRequest,
    IngestRequest,
    ReadStatusUpdate,
    StatsResponse,
    UploadResponse,
)

STORAGE_ROOT = Path(os.getenv("HUB_BASE_DIR", str(Path.home() / "devops-knowledge-hub"))) / "documents"

app = FastAPI(title="DevOps Knowledge Hub API", version="2.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_origin_regex=r"chrome-extension://.*",
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# In-memory batch progress (single-process local tool)
_batch: dict = {
    "status": "idle", "total": 0, "processed": 0,
    "failed": [], "current_title": None, "current_domain": None,
    "done": True, "task_id": None,
}


@app.on_event("startup")
async def startup():
    STORAGE_ROOT.mkdir(parents=True, exist_ok=True)
    database.init_db()
    log.info("DevOps Knowledge Hub v2.0 started")

    # Start watch folder in daemon thread
    watch_dir = os.getenv("WATCH_DIR", "")
    if watch_dir:
        loop = asyncio.get_event_loop()
        from watcher import start_watcher
        t = threading.Thread(target=start_watcher, args=(loop,), daemon=True)
        t.start()


# ── Document upload ──────────────────────────────────────────────────────────

@app.post("/api/documents/upload", response_model=UploadResponse)
async def upload_document(file: UploadFile = File(...), metadata: str = Form("{}")):
    raw_meta = json.loads(metadata)
    content = await file.read()
    file_hash = hashlib.sha256(content).hexdigest()

    existing = database.find_by_hash(file_hash)
    if existing:
        return UploadResponse(id=existing["id"], status="duplicate", path=existing["file_path"])

    original_filename = file.filename or "document"
    suffix = Path(original_filename).suffix.lower() or _guess_suffix(file.content_type)
    dest_path = STORAGE_ROOT / f"{file_hash}{suffix}"
    dest_path.write_bytes(content)

    file_type = suffix.lstrip(".")
    doc_id = str(uuid.uuid4())
    doc = {
        "id": doc_id, "hash": file_hash,
        "title": raw_meta.get("title") or Path(original_filename).stem,
        "filename": original_filename,
        "file_path": str(dest_path), "file_type": file_type,
        "file_size": len(content),
        "source_url": raw_meta.get("source_url"),
        "author": raw_meta.get("author"),
        "linkedin_post_url": raw_meta.get("linkedin_post_url"),
        "domain": "uncategorized", "tags": [], "summary": None,
        "read_status": "unread",
        "raw_text": raw_meta.get("raw_text"),
        "doc_type": "raw", "ingest_source": "linkedin_extension",
    }
    database.insert_document(doc)
    asyncio.create_task(categorize_document(doc_id))
    return UploadResponse(id=doc_id, status="saved", path=str(dest_path))


@app.post("/api/documents/upload-images", response_model=UploadResponse)
async def upload_images(body: ImageUploadRequest):
    if not body.images:
        raise HTTPException(400, "No images provided")

    pdf_doc = fitz.open()
    for data_url in body.images:
        try:
            if "," in data_url:
                header, b64 = data_url.split(",", 1)
                filetype = "png" if "image/png" in header else "jpeg"
            else:
                b64, filetype = data_url, "jpeg"
            img_bytes = base64.b64decode(b64)
            img_doc = fitz.open(stream=img_bytes, filetype=filetype)
            page_pdf = fitz.open("pdf", img_doc.convert_to_pdf())
            pdf_doc.insert_pdf(page_pdf)
            img_doc.close(); page_pdf.close()
        except Exception:
            continue

    if pdf_doc.page_count == 0:
        raise HTTPException(422, "No valid images could be processed")

    raw_pdf = pdf_doc.tobytes(garbage=4, deflate=True)
    pdf_doc.close()

    file_hash = hashlib.sha256(raw_pdf).hexdigest()
    existing = database.find_by_hash(file_hash)
    if existing:
        return UploadResponse(id=existing["id"], status="duplicate", path=existing["file_path"])

    dest_path = STORAGE_ROOT / f"{file_hash}.pdf"
    dest_path.write_bytes(raw_pdf)
    meta = body.metadata
    doc_id = str(uuid.uuid4())
    doc = {
        "id": doc_id, "hash": file_hash,
        "title": (meta.title if meta else None) or "Captured LinkedIn Document",
        "filename": f"{file_hash}.pdf",
        "file_path": str(dest_path), "file_type": "pdf",
        "file_size": len(raw_pdf),
        "source_url": meta.source_url if meta else None,
        "author": meta.author if meta else None,
        "linkedin_post_url": meta.linkedin_post_url if meta else None,
        "domain": "uncategorized", "tags": [], "summary": None,
        "read_status": "unread",
        "raw_text": meta.post_text if meta else None,
        "doc_type": "raw", "ingest_source": "linkedin_extension",
    }
    database.insert_document(doc)
    asyncio.create_task(categorize_document(doc_id))
    return UploadResponse(id=doc_id, status="saved", path=str(dest_path))


# ── Document CRUD ────────────────────────────────────────────────────────────

@app.get("/api/documents", response_model=list[DocumentListItem])
def list_documents(
    domain: Optional[str] = None,
    search: Optional[str] = None,
    limit: int = 100,
    offset: int = 0,
):
    rows = database.list_documents(domain=domain, search=search, limit=limit, offset=offset)
    return [_to_list_item(r) for r in rows]


@app.get("/api/documents/{doc_id}", response_model=DocumentRecord)
def get_document(doc_id: str):
    doc = database.get_document(doc_id)
    if not doc:
        raise HTTPException(404, "Document not found")
    return doc


@app.get("/api/documents/{doc_id}/file")
def serve_file(doc_id: str):
    doc = database.get_document(doc_id)
    if not doc:
        raise HTTPException(404, "Document not found")
    file_path = Path(doc["file_path"])
    if not file_path.exists():
        raise HTTPException(404, "File not found on disk")
    media_type, _ = mimetypes.guess_type(str(file_path))
    return FileResponse(
        path=str(file_path),
        media_type=media_type or "application/octet-stream",
        filename=doc["filename"],
    )


@app.get("/api/documents/{doc_id}/view")
def view_file(doc_id: str):
    """Serve file inline for browser rendering — no download prompt."""
    doc = database.get_document(doc_id)
    if not doc:
        raise HTTPException(404, "Document not found")
    file_path = Path(doc["file_path"])
    if not file_path.exists():
        raise HTTPException(404, "File not found on disk")

    file_type = doc.get("file_type", "").lower()

    if file_type == "pdf":
        return Response(
            content=file_path.read_bytes(),
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'inline; filename="{doc["filename"]}"',
                "X-Frame-Options": "ALLOWALL",
                "Cache-Control": "private, max-age=3600",
            },
        )
    if file_type in ("pptx", "ppt"):
        return HTMLResponse(_pptx_to_html(str(file_path), doc))
    if file_type in ("docx", "doc"):
        return HTMLResponse(_docx_to_html(str(file_path), doc))
    if file_type in ("md", "markdown", "txt"):
        text = file_path.read_text(errors="replace")
        return HTMLResponse(_md_to_html(text, doc))

    media_type, _ = mimetypes.guess_type(str(file_path))
    return FileResponse(
        path=str(file_path),
        media_type=media_type or "application/octet-stream",
        filename=doc["filename"],
    )


@app.get("/api/documents/{doc_id}/concept-card")
def get_concept_card(doc_id: str):
    card = database.get_concept_card_for_doc(doc_id)
    if not card:
        raise HTTPException(404, "No concept card for this document")
    return card


@app.get("/api/documents/{doc_id}/similar")
def get_similar_docs(doc_id: str):
    return database.get_similar_docs(doc_id)


@app.patch("/api/documents/{doc_id}/domain")
def update_domain(doc_id: str, body: DomainUpdate):
    if not database.get_document(doc_id):
        raise HTTPException(404, "Document not found")
    body.validate_domain()
    return database.update_document(doc_id, {"domain": body.domain})


@app.patch("/api/documents/{doc_id}/status")
def update_read_status(doc_id: str, body: ReadStatusUpdate):
    if not database.get_document(doc_id):
        raise HTTPException(404, "Document not found")
    return database.update_document(doc_id, {"read_status": body.read_status})


# ── Categorization ───────────────────────────────────────────────────────────

@app.post("/api/documents/{doc_id}/categorize")
async def trigger_categorize(doc_id: str):
    if not database.get_document(doc_id):
        raise HTTPException(404, "Document not found")
    updated = await recategorize_document(doc_id)
    return updated or database.get_document(doc_id)


@app.post("/api/categorize/batch")
async def start_batch_categorize(body: BatchCategorizeRequest):
    global _batch
    if _batch["status"] == "running":
        return {"status": "already_running", "task_id": _batch["task_id"]}

    task_id = str(uuid.uuid4())
    _batch.update({
        "status": "starting", "total": 0, "processed": 0,
        "failed": [], "current_title": None, "current_domain": None,
        "done": False, "task_id": task_id,
    })
    asyncio.create_task(_run_batch(body.limit, body.only_uncategorized))
    return {"status": "started", "task_id": task_id}


@app.post("/api/categorize/{doc_id}")
async def categorize_single(doc_id: str):
    if not database.get_document(doc_id):
        raise HTTPException(404, "Document not found")
    updated = await recategorize_document(doc_id)
    return updated or database.get_document(doc_id)


async def _run_batch(limit: int, only_uncategorized: bool):
    global _batch
    _batch["status"] = "running"

    async def on_progress(processed, total, current_title, current_domain):
        _batch.update({
            "processed": processed, "total": total,
            "current_title": current_title, "current_domain": current_domain,
        })

    try:
        result = await batch_categorize(
            limit=limit,
            only_uncategorized=only_uncategorized,
            progress_callback=on_progress,
        )
        _batch.update({
            "status": "done", "processed": result["processed"],
            "total": result["total"], "failed": result["failed"], "done": True,
        })
    except Exception as e:
        log.error("Batch error: %s", e)
        _batch["status"] = "error"
        _batch["done"] = True


@app.get("/api/categorize/progress")
async def stream_batch_progress():
    async def event_gen() -> AsyncGenerator[str, None]:
        last = -1
        while True:
            data = json.dumps({
                "status": _batch["status"],
                "total": _batch["total"],
                "processed": _batch["processed"],
                "failed": len(_batch["failed"]),
                "current_title": _batch["current_title"],
                "current_domain": _batch["current_domain"],
                "done": _batch["done"],
            })
            yield f"data: {data}\n\n"
            if _batch["done"]:
                break
            await asyncio.sleep(1.0)

    return StreamingResponse(
        event_gen(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


# ── Concept extraction ────────────────────────────────────────────────────────

@app.post("/api/documents/{doc_id}/extract-concepts")
async def extract_concepts_endpoint(doc_id: str):
    if not database.get_document(doc_id):
        raise HTTPException(404, "Document not found")
    from extractor import extract_concepts
    result = await extract_concepts(doc_id)
    if not result:
        raise HTTPException(503, "Concept extraction failed — check GEMINI_API_KEY")
    return result


# ── Search & Q&A ──────────────────────────────────────────────────────────────

@app.post("/api/search")
async def search_documents(body: dict):
    from search import search
    query = body.get("query", "")
    if not query:
        raise HTTPException(400, "query is required")
    return await search(query=query, limit=body.get("limit", 10), domain=body.get("domain"))


@app.post("/api/ask")
async def ask_question(body: dict):
    from search import ask
    question = body.get("question", "")
    if not question:
        raise HTTPException(400, "question is required")
    return await ask(
        question=question,
        doc_ids=body.get("doc_ids", []),
        domain=body.get("domain"),
        history=body.get("history", []),
    )


# ── Chat ──────────────────────────────────────────────────────────────────────

@app.post("/api/chat")
async def chat_endpoint(body: dict):
    import ai_client
    message = body.get("message", "")
    if not message:
        raise HTTPException(400, "message is required")

    # accept both doc_id (single) and doc_ids (array from portal chat widget)
    doc_id = body.get("doc_id") or (body.get("doc_ids") or [None])[0]
    domain = body.get("domain")
    history = body.get("history", [])
    context, sources = "", []

    if doc_id:
        card = database.get_concept_card_for_doc(str(doc_id))
        doc = card or database.get_document(str(doc_id))
        if doc:
            context = (doc.get("raw_text") or doc.get("summary") or "")[:4000]
            sources.append(doc.get("title") or doc.get("filename"))
    elif domain:
        guide = database.get_study_guide(domain)
        if guide:
            context = guide["content"][:6000]
            sources.append(f"{domain} study guide")

    system = (
        "You are an expert SRE/DevOps assistant for a Staff/Principal engineer. "
        "Give direct, technically deep answers. No installation steps. "
        "Focus on concepts, trade-offs, production patterns, and interview-relevant depth."
    )
    full_prompt = f"{system}\n\n"
    if context:
        full_prompt += f"Context from documents:\n{context}\n\n"
    full_prompt += f"Question: {message}"

    try:
        chat = ai_client.get_chat_session(history, task="chat")
        response = chat.send_message(full_prompt)
        return {"response": response.text, "sources": sources}
    except Exception as e:
        raise HTTPException(503, f"Chat error: {e}")


# ── Study guides ──────────────────────────────────────────────────────────────

@app.post("/api/study-guide/generate")
async def generate_study_guide_endpoint(body: dict):
    from study_guide import generate_study_guide
    domain = body.get("domain")
    if not domain:
        raise HTTPException(400, "domain is required")
    result = await generate_study_guide(domain, min_docs=1)
    if not result:
        raise HTTPException(503, "Study guide generation failed")
    return result


@app.get("/api/study-guide/{domain}")
def get_study_guide(domain: str):
    guide = database.get_study_guide(domain)
    if not guide:
        raise HTTPException(404, "Study guide not found for this domain")
    return guide


_guide_gen_status: dict = {"running": False, "done": 0, "total": 0, "current": "", "errors": []}


@app.put("/api/study-guide/{domain}")
async def save_study_guide_directly(domain: str, body: dict):
    """Save guide content directly — used when content is generated outside Gemini."""
    content = body.get("content", "").strip()
    if not content:
        raise HTTPException(400, "content is required")
    doc_count = body.get("doc_count", 0)
    from pathlib import Path
    guides_dir = Path(__file__).parent.parent / "study-guides"
    guides_dir.mkdir(parents=True, exist_ok=True)
    (guides_dir / f"{domain}.md").write_text(content, encoding="utf-8")
    database.save_study_guide(domain, content, doc_count)
    return {"domain": domain, "chars": len(content), "saved": True}


@app.post("/api/study-guide/generate-all")
async def generate_all_study_guides_endpoint(background_tasks: BackgroundTasks):
    global _guide_gen_status
    if _guide_gen_status.get("running"):
        return {"status": "already_running", **_guide_gen_status}
    _guide_gen_status = {"running": True, "done": 0, "total": 0, "current": "", "errors": []}

    async def _run():
        global _guide_gen_status
        from study_guide import generate_all_study_guides
        try:
            results = await generate_all_study_guides(min_docs=3, status=_guide_gen_status)
        except Exception as e:
            _guide_gen_status["errors"].append(str(e))
        finally:
            _guide_gen_status["running"] = False

    asyncio.create_task(_run())
    return {"status": "started", "message": "Study guide generation running in background. Poll /api/study-guide/generate-all/status"}


@app.get("/api/study-guide/generate-all/status")
def guide_gen_status():
    return _guide_gen_status


@app.post("/api/study-guide/{domain}/patch")
async def patch_study_guide_endpoint(domain: str, body: dict):
    """Smart append — reads new doc, adds ONLY genuinely new content to existing guide."""
    from study_guide import patch_study_guide_with_doc, generate_study_guide
    doc_id = body.get("doc_id")

    if not doc_id:
        return {"error": "doc_id required"}
    if not database.get_document(doc_id):
        raise HTTPException(404, "Document not found")

    result = await patch_study_guide_with_doc(domain, doc_id)
    return result or {"error": "Patch failed — check backend logs"}


# ── Ingestion (Phase 6) ───────────────────────────────────────────────────────

@app.post("/api/ingest")
async def ingest_endpoint(body: IngestRequest):
    from ingestion import ingest
    result = await ingest(
        source_type=body.source_type,
        source_data=body.source_data,
        options=body.options or {},
    )
    return result


@app.post("/api/ingest/bulk")
async def ingest_bulk(body: dict):
    from ingestion import ingest
    items = body.get("items", [])
    if not items:
        raise HTTPException(400, "items list is required")
    results = []
    for item in items:
        try:
            r = await ingest(
                source_type=item.get("source_type", "url"),
                source_data=item.get("source_data", ""),
                options=item.get("options", {}),
            )
            results.append(r)
        except Exception as e:
            results.append({"status": "error", "message": str(e), "source": item.get("source_data")})
    return {"total": len(items), "results": results}


# ── Discovery (Phase 6) ───────────────────────────────────────────────────────

@app.post("/api/discover")
async def discover_endpoint(body: DiscoverRequest):
    from discovery import discover
    return await discover(
        query=body.query,
        mode=body.mode,
        reference_doc_id=body.reference_doc_id,
    )


@app.post("/api/discover/download")
async def discover_download(body: DiscoverDownloadRequest):
    from discovery import download_selected
    results = await download_selected(body.session_id, body.selected_urls)
    return {"results": results, "total": len(results)}


@app.get("/api/discover/sessions")
def list_discover_sessions():
    return database.list_discovery_sessions(limit=20)


@app.get("/api/discover/status/{session_id}")
async def stream_discover_status(session_id: str):
    async def event_gen() -> AsyncGenerator[str, None]:
        for _ in range(60):
            session = database.get_discovery_session(session_id)
            if not session:
                yield f"data: {json.dumps({'status': 'not_found'})}\n\n"
                break
            yield f"data: {json.dumps({'status': session['status'], 'session_id': session_id})}\n\n"
            if session["status"] in ("complete", "failed"):
                break
            await asyncio.sleep(2.0)

    return StreamingResponse(
        event_gen(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.get("/api/watch/events")
def watch_events():
    return database.list_watch_events(limit=50)


# ── Admin ─────────────────────────────────────────────────────────────────────

@app.post("/api/admin/generate-portal")
def generate_portal():
    import subprocess, sys
    script = Path(__file__).parent / "generate_portal.py"
    try:
        result = subprocess.run(
            [sys.executable, str(script)],
            capture_output=True, text=True, timeout=180,
        )
        return {
            "success": result.returncode == 0,
            "output": result.stdout,
            "error": result.stderr if result.returncode != 0 else None,
        }
    except Exception as e:
        raise HTTPException(500, str(e))


# ── Stats & health ────────────────────────────────────────────────────────────

@app.get("/api/stats")
def get_stats():
    return database.get_stats()


@app.get("/api/health")
def health():
    import ai_client
    return {
        "status": "ok",
        "ai_configured": ai_client.is_configured(),
        "batch_status": _batch["status"],
        "batch_processed": _batch["processed"],
        "batch_total": _batch["total"],
    }


# ── Helpers ───────────────────────────────────────────────────────────────────

def _to_list_item(row: dict) -> DocumentListItem:
    return DocumentListItem(
        id=row["id"], title=row.get("title"), filename=row["filename"],
        file_type=row["file_type"], file_size=row["file_size"],
        author=row.get("author"), domain=row["domain"],
        tags=row.get("tags", []), summary=row.get("summary"),
        read_status=row["read_status"], saved_at=row["saved_at"],
        linkedin_post_url=row.get("linkedin_post_url"),
        source_url=row.get("source_url"),
        concept_level=row.get("concept_level"),
        key_topics=row.get("key_topics"),
        doc_type=row.get("doc_type", "raw"),
    )


def _guess_suffix(content_type: Optional[str]) -> str:
    if not content_type:
        return ".bin"
    return {
        "application/pdf": ".pdf",
        "application/vnd.ms-powerpoint": ".ppt",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
        "application/vnd.ms-excel": ".xls",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
        "application/msword": ".doc",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    }.get(content_type, ".bin")


def _pptx_to_html(file_path: str, doc: dict) -> str:
    body = "<p>PPTX preview requires python-pptx.</p>"
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                s.text.strip() for s in slide.shapes
                if hasattr(s, "text") and s.text.strip()
            ]
            slides.append(
                f'<div class="slide"><strong>Slide {i}</strong>'
                f'<p>{"<br>".join(_html.escape(t) for t in texts)}</p></div>'
            )
        body = "\n".join(slides) or "<p>No text content found.</p>"
    except Exception:
        pass
    return _wrap_html(doc.get("title") or doc.get("filename", "Document"), body)


def _docx_to_html(file_path: str, doc: dict) -> str:
    body = "<p>DOCX preview requires python-docx.</p>"
    try:
        from docx import Document
        d = Document(file_path)
        body = "\n".join(
            f"<p>{_html.escape(p.text)}</p>" for p in d.paragraphs if p.text.strip()
        )
    except Exception:
        pass
    return _wrap_html(doc.get("title") or doc.get("filename", "Document"), body)


def _md_to_html(text: str, doc: dict) -> str:
    lines = text.split("\n")
    parts = []
    for line in lines:
        if line.startswith("# "):
            parts.append(f"<h1>{_html.escape(line[2:])}</h1>")
        elif line.startswith("## "):
            parts.append(f"<h2>{_html.escape(line[3:])}</h2>")
        elif line.startswith("### "):
            parts.append(f"<h3>{_html.escape(line[4:])}</h3>")
        elif line.startswith(("- ", "* ")):
            parts.append(f"<li>{_html.escape(line[2:])}</li>")
        elif line.startswith("> "):
            parts.append(f"<blockquote>{_html.escape(line[2:])}</blockquote>")
        elif line.strip() == "":
            parts.append("<br>")
        else:
            parts.append(f"<p>{_html.escape(line)}</p>")
    return _wrap_html(doc.get("title") or doc.get("filename", "Document"), "\n".join(parts))


def _wrap_html(title: str, body: str) -> str:
    return f"""<!DOCTYPE html><html lang="en"><head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>{_html.escape(title)}</title>
<style>
body{{font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;
     max-width:900px;margin:0 auto;padding:32px 24px;line-height:1.7;
     color:#1a1a2e;background:#fff}}
h1{{font-size:1.8em;border-bottom:2px solid #e0e0e0;padding-bottom:12px}}
h2{{font-size:1.3em;color:#0a66c2;margin-top:2em}}
h3{{font-size:1.1em}}
.slide{{border-left:3px solid #0a66c2;padding:8px 16px;margin:12px 0}}
blockquote{{border-left:3px solid #ddd;padding-left:16px;color:#666}}
li{{margin:4px 0}}
</style></head><body>
<h1>{_html.escape(title)}</h1>
{body}
</body></html>"""


if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=int(os.getenv("BACKEND_PORT", "8765")), reload=True)
